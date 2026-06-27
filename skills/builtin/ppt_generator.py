"""PPT 生成技能 — python-pptx"""
from pathlib import Path
from api.agent_llm import call_llm

skill_def = {
    "name": "ppt-generator", "version": "1.0.0",
    "description": "PPT 演示文稿生成",
    "author": "AUTO-EVO-AI", "category": "文件生成", "icon": "📊",
    "tags": ["PPT", "演示", "幻灯片"],
    "input_schema": {"type": "object", "properties": {"topic": {"type": "string"}, "pages": {"type": "integer"}}},
    "output_schema": {"type": "object", "properties": {"file_path": {"type": "string"}}}
}

OUT = Path(__file__).resolve().parent.parent.parent / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)

def execute(params, context=None):
    topic = params.get("topic", "")
    pages = int(params.get("pages", 5))
    if not topic:
        return {"file_path": "", "error": "请提供 PPT 主题（topic）"}
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        # 用 LLM 生成每页内容
        sp = f"请为'{topic}'生成一个{pages}页的PPT大纲，每页格式：标题: 内容要点（用逗号分隔），直接输出不要额外解释。"
        text, _ = call_llm([{"role": "user", "content": sp}])
        lines = [l for l in (text or "").split("\n") if ":" in l]
        if not lines:
            lines = [f"{topic}: 介绍, 核心功能, 应用场景, 优势, 总结"]
        for line in lines[:pages]:
            if ":" not in line: continue
            title, rest = line.split(":", 1)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title.strip()
            items = [x.strip() for x in rest.split(",") if x.strip()]
            tf = slide.placeholders[1].text_frame
            for i, item in enumerate(items):
                if i == 0:
                    tf.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
        fp = str(OUT / f"{topic.replace(' ', '_')[:20]}.pptx")
        prs.save(fp)
        return {"file_path": fp}
    except ImportError:
        return {"file_path": "", "error": "python-pptx 未安装，请执行 pip install python-pptx"}
    except Exception as e:
        return {"file_path": "", "error": f"PPT 生成失败：{e}"}
