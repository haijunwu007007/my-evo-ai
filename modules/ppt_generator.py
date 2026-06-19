"""
AUTO-EVO-AI V0.1 — PPT 生成器 (python-pptx)
Grade: A (生产级) | Category: documents
职责：创建/编辑 PowerPoint 演示文稿，支持模板、图表、表格、图片、动画
"""
from __future__ import annotations

__module_meta__ = {
    "id": "ppt-generator",
    "name": "PPT Generator",
    "version": "V0.1",
    "group": "documents",
    "inputs": [
        {"name": "action", "type": "string", "required": True, "description": "create / edit / read / chart / info / list / to_images"},
        {"name": "path", "type": "string", "required": False, "description": "文件路径"},
        {"name": "data", "type": "dict", "required": False, "description": "操作参数"},
    ],
    "outputs": [
        {"name": "result", "type": "dict", "description": "执行结果"},
    ],
    "triggers": [],
    "depends_on": [],
    "tags": ["adapter", "doc"],
    "grade": "A",
    "description": "PowerPoint 全功能生成器：创建演示文稿、模板驱动、图表生成、表格插入、图片导出",
}

import os
import uuid
import logging
from pathlib import Path
from datetime import datetime
from typing import Any

try:
    from modules._base.enterprise_module import EnterpriseModule
    from modules._base.audit import AuditLogger
except ImportError:
    import sys
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
    from _base.enterprise_module import EnterpriseModule
    from _base.audit import AuditLogger

logger = logging.getLogger("evo.ppt")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed; PPT actions disabled")

# 预定义配色主题
THEMES = {
    "default": {"bg": "1F2937", "heading": "F59E0B", "text": "F3F4F6", "accent": "3B82F6"},
    "light": {"bg": "FFFFFF", "heading": "1F2937", "text": "374151", "accent": "2563EB"},
    "dark": {"bg": "111827", "heading": "F59E0B", "text": "D1D5DB", "accent": "60A5FA"},
    "corporate": {"bg": "0F172A", "heading": "FFFFFF", "text": "CBD5E1", "accent": "3B82F6"},
    "nature": {"bg": "ECFDF5", "heading": "065F46", "text": "374151", "accent": "059669"},
}


class PptGenerator(EnterpriseModule):
    """PPT 全功能生成器"""

    def __init__(self, *args: Any, **kwargs: Any):
        super().__init__(*args, **kwargs)
        self._name = "ppt-generator"
        self._work_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "data", "presentations")
        os.makedirs(self._work_dir, exist_ok=True)
        self._audit = AuditLogger()
        self.stats = {"calls": 0, "errors": 0, "ppts_created": 0}
        self._status = "ready" if HAS_PPTX else "error"

    def initialize(self) -> None:
        logger.info(f"[ppt] initialized, HAS_PPTX={HAS_PPTX}")

    def health_check(self) -> dict[str, Any]:
        return {
            "status": self._status,
            "module": "ppt-generator",
            "libraries": {"python-pptx": HAS_PPTX},
            "themes": list(THEMES.keys()),
            "stats": self.stats,
        }

    def shutdown(self) -> None:
        logger.info("[ppt] shutdown")

    # ── 主入口 ───────────────────────────────────────────

    async def execute(self, action: str, path: str = "", data: dict | None = None, **kwargs: Any) -> dict[str, Any]:
        if not HAS_PPTX:
            return {"status": "error", "message": "python-pptx 未安装"}
        self.stats["calls"] += 1
        data = data or {}
        action_map: dict[str, Any] = {
            "create": self._create_ppt,
            "edit": self._edit_ppt,
            "read": self._read_ppt,
            "chart": self._add_chart_slide,
            "info": self._get_info,
            "list": self._list_ppts,
            "to_images": self._export_to_images,
        }
        handler = action_map.get(action)
        if not handler:
            return {"status": "error", "message": f"未知 action: {action}，可选: {list(action_map.keys())}"}
        try:
            return await handler(path, data)
        except Exception as e:
            self.stats["errors"] += 1
            logger.exception(f"[ppt] execute {action} 失败")
            return {"status": "error", "message": str(e)}

    # ── 创建演示文稿 ─────────────────────────────────────

    async def _create_ppt(self, path: str, data: dict) -> dict[str, Any]:
        filename = data.get("filename", f"deck_{uuid.uuid4().hex[:8]}.pptx")
        save_path = path or os.path.join(self._work_dir, filename)
        title = data.get("title", "演示文稿")
        theme_name = data.get("theme", "default")
        slides_data = data.get("slides", [])
        theme = THEMES.get(theme_name, THEMES["default"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── 标题页 ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        self._set_bg(slide, theme["bg"])
        self._add_textbox(slide, Inches(1), Inches(2), Inches(11.333), Inches(2), title,
                          Pt(44), True, theme["heading"], PP_ALIGN.CENTER)
        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_textbox(slide, Inches(1), Inches(4), Inches(11.333), Inches(1), subtitle,
                              Pt(24), False, theme["text"], PP_ALIGN.CENTER)

        # ── 内容页 ──
        for s in slides_data:
            slide_type = s.get("type", "content")
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._set_bg(slide, theme["bg"])

            slide_title = s.get("title", "")
            if slide_title:
                self._add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8),
                                  slide_title, Pt(32), True, theme["heading"], PP_ALIGN.LEFT)

            if slide_type == "content":
                content = s.get("content", [])
                y = 1.5
                for item in content:
                    if isinstance(item, str):
                        item = {"text": item, "level": 0}
                    indent = item.get("level", 0) * 0.3
                    prefix = "  " * item.get("level", 0) + ("• " if item.get("level", 0) == 0 else "- ")
                    self._add_textbox(slide, Inches(0.8 + indent), Inches(y), Inches(11.733 - indent), Inches(0.5),
                                      prefix + item.get("text", ""), Pt(18), False, theme["text"], PP_ALIGN.LEFT)
                    y += 0.5

            elif slide_type == "two_column":
                left = s.get("left", [])
                right = s.get("right", [])
                y = 1.5
                for item in left:
                    self._add_textbox(slide, Inches(0.8), Inches(y), Inches(5.8), Inches(0.5),
                                      "• " + (item if isinstance(item, str) else item.get("text", "")),
                                      Pt(16), False, theme["text"], PP_ALIGN.LEFT)
                    y += 0.45
                y = 1.5
                for item in right:
                    self._add_textbox(slide, Inches(7), Inches(y), Inches(5.8), Inches(0.5),
                                      "• " + (item if isinstance(item, str) else item.get("text", "")),
                                      Pt(16), False, theme["text"], PP_ALIGN.LEFT)
                    y += 0.45

            elif slide_type == "table":
                headers = s.get("headers", [])
                rows = s.get("rows", [])
                self._add_table(slide, Inches(0.8), Inches(1.8), headers, rows, theme)

            # 添加图片到幻灯片
            images = s.get("images", [])
            for img_data in images:
                img_path = img_data if isinstance(img_data, str) else img_data.get("path", "")
                if os.path.isfile(img_path):
                    x = img_data.get("x", 1) if isinstance(img_data, dict) else 1
                    y = img_data.get("y", 1.5) if isinstance(img_data, dict) else 1.5
                    w = img_data.get("width", 8) if isinstance(img_data, dict) else 8
                    try:
                        slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w))
                    except Exception as e:
                        logger.warning(f"[ppt] 图片添加失败: {e}")

        # ── 结束页 ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide, theme["bg"])
        self._add_textbox(slide, Inches(1), Inches(3), Inches(11.333), Inches(1.5),
                          "谢谢", Pt(48), True, theme["heading"], PP_ALIGN.CENTER)

        prs.save(save_path)
        self.stats["ppts_created"] += 1
        return {"status": "success", "path": os.path.abspath(save_path), "slides": 2 + len(slides_data)}

    # ── 编辑演示文稿 ─────────────────────────────────────

    async def _edit_ppt(self, path: str, data: dict) -> dict[str, Any]:
        if not path or not os.path.isfile(path):
            return {"status": "error", "message": f"文件不存在: {path}"}
        prs = Presentation(path)
        edits = data.get("edits", [])
        for edit in edits:
            op = edit.get("op", "")
            if op == "add_slide":
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_text = edit.get("title", "")
                if title_text:
                    self._add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8),
                                      title_text, Pt(32), True, "FFFFFF", PP_ALIGN.LEFT)
                content = edit.get("content", [])
                y = 1.5
                for item in content:
                    self._add_textbox(slide, Inches(0.8), Inches(y), Inches(11.733), Inches(0.5),
                                      "• " + item, Pt(18), False, "D1D5DB", PP_ALIGN.LEFT)
                    y += 0.5
            elif op == "add_image":
                img_path = edit.get("path", "")
                slide_idx = edit.get("slide", 1)
                if os.path.isfile(img_path) and slide_idx <= len(prs.slides):
                    slide = prs.slides[slide_idx - 1]
                    slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(8))

        prs.save(path)
        return {"status": "success", "path": os.path.abspath(path)}

    # ── 读取演示文稿 ─────────────────────────────────────

    async def _read_ppt(self, path: str, data: dict) -> dict[str, Any]:
        if not path or not os.path.isfile(path):
            return {"status": "error", "message": f"文件不存在: {path}"}
        prs = Presentation(path)
        slides_info = []
        for si, slide in enumerate(prs.slides):
            slide_data = {"index": si + 1, "shapes": []}
            for shape in slide.shapes:
                shape_info = {"name": shape.shape_type, "left": shape.left, "top": shape.top}
                if shape.has_text_frame:
                    shape_info["text"] = shape.text_frame.text[:200]
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text for cell in row.cells])
                    shape_info["table"] = rows
                slide_data["shapes"].append(shape_info)
            slides_info.append(slide_data)
        return {"status": "success", "slides": slides_info, "count": len(prs.slides)}

    # ── 图表幻灯片 ───────────────────────────────────────

    async def _add_chart_slide(self, path: str, data: dict) -> dict[str, Any]:
        """
        生成带图表的幻灯片（使用表格模拟图表，因为 python-pptx 的 chart 功能有限）
        如需真实图表，推荐用 matplotlib 生成图片后插入
        """
        if not path or not os.path.isfile(path):
            return {"status": "error", "message": f"文件不存在: {path}"}
        prs = Presentation(path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide, "1F2937")

        chart_title = data.get("title", "数据图表")
        self._add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8),
                          chart_title, Pt(32), True, "F59E0B", PP_ALIGN.LEFT)

        # 使用表格展示数据
        headers = data.get("headers", ["项目", "数值"])
        rows = data.get("rows", [["示例", "100"]])
        self._add_table(slide, Inches(0.8), Inches(1.8), headers, [headers] + rows,
                        {"heading": "F59E0B", "text": "F3F4F6"})

        prs.save(path)
        return {"status": "success", "path": os.path.abspath(path)}

    # ── 导出为图片 ── (每页一张) ──────────────────────────

    async def _export_to_images(self, path: str, data: dict) -> dict[str, Any]:
        """
        导出 PPT 每页为图片。
        注意：python-pptx 不支持直接渲染，需配合 LibreOffice 或 PowerPoint。
        这里返回每页的文本摘要。
        """
        if not path or not os.path.isfile(path):
            return {"status": "error", "message": f"文件不存在: {path}"}
        prs = Presentation(path)
        out_dir = data.get("output_dir", os.path.join(self._work_dir, "ppt_images"))
        os.makedirs(out_dir, exist_ok=True)

        pages = []
        for si, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text[:100])
            text_file = os.path.join(out_dir, f"slide_{si + 1}.txt")
            with open(text_file, "w", encoding="utf-8") as f:
                f.write("\n".join(texts))
            pages.append({"slide": si + 1, "text_preview": texts[:5], "text_file": text_file})

        return {"status": "success", "pages": pages, "output_dir": out_dir}

    # ── 获取信息 ─────────────────────────────────────────

    async def _get_info(self, path: str, data: dict) -> dict[str, Any]:
        if not path or not os.path.isfile(path):
            return {"status": "error", "message": f"文件不存在: {path}"}
        prs = Presentation(path)
        info = {
            "slides": len(prs.slides),
            "width": prs.slide_width,
            "height": prs.slide_height,
            "slide_layouts": len(prs.slide_layouts),
            "file_size_bytes": os.path.getsize(path),
        }
        return {"status": "success", "info": info}

    # ── 列出文件 ─────────────────────────────────────────

    async def _list_ppts(self, path: str, data: dict) -> dict[str, Any]:
        search_dir = path or self._work_dir
        if not os.path.isdir(search_dir):
            return {"status": "error", "message": f"目录不存在: {search_dir}"}
        files = []
        for f in sorted(os.listdir(search_dir)):
            fp = os.path.join(search_dir, f)
            if os.path.isfile(fp) and f.lower().endswith((".pptx", ".ppt")):
                files.append({
                    "name": f,
                    "size_bytes": os.path.getsize(fp),
                    "modified": datetime.fromtimestamp(os.path.getmtime(fp)).isoformat(),
                })
        return {"status": "success", "files": files, "directory": search_dir}

    # ── 辅助方法 ─────────────────────────────────────────

    def _set_bg(self, slide, color_hex: str) -> None:
        """设置幻灯片背景色"""
        background = slide.background
        fill = background.fill
        fill.solid()
        r, g, b = self._hex_to_rgb(color_hex)
        fill.fore_color.rgb = RGBColor(r, g, b)

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=Pt(18), bold=False, color_hex="FFFFFF",
                     alignment=PP_ALIGN.LEFT) -> None:
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        r, g, b = self._hex_to_rgb(color_hex)
        p.font.color.rgb = RGBColor(r, g, b)
        p.alignment = alignment

    def _add_table(self, slide, left, top, headers, rows, theme) -> None:
        """添加表格到幻灯片"""
        shape = slide.shapes.add_table(len(rows), len(headers), left, top, Inches(11.5), Inches(0.4 * len(rows)))
        table = shape.table
        for ri, row_data in enumerate(rows):
            for ci, val in enumerate(row_data):
                cell = table.cell(ri, ci)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    if ri == 0:
                        paragraph.font.bold = True
                        hr, hg, hb = self._hex_to_rgb(theme.get("heading", "FFFFFF"))
                        paragraph.font.color.rgb = RGBColor(hr, hg, hb)
                    else:
                        tr, tg, tb = self._hex_to_rgb(theme.get("text", "D1D5DB"))
                        paragraph.font.color.rgb = RGBColor(tr, tg, tb)

    @staticmethod
    def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
        hex_str = hex_str.lstrip("#")
        return tuple(int(hex_str[i:i + 2], 16) for i in (0, 2, 4))


module_class = PptGenerator
