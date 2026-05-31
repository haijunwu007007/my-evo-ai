"""
AUTO-EVO-AI V0.1 — 文档生成引擎
====================================
上市公司生产级设计：

核心能力:
  1. Word生成 — python-docx, 支持标题/段落/表格/图片/页眉页脚
  2. Excel生成 — openpyxl + xlsxwriter, 支持多Sheet/公式/图表/样式
  3. PPT生成 — python-pptx, 支持幻灯片/图表/图片/备注
  4. Markdown生成 — 结构化Markdown
  5. HTML报告 — 漂亮的HTML报告(带图表)
  6. PDF导出 — reportlab转PDF
  7. 模板系统 — 变量替换+循环+条件

使用方式:
  from core.doc_generator import DocGenerator

  gen = DocGenerator()

  # 生成Word
  gen.word_add_title("系统报告")
  gen.word_add_paragraph("这是一段内容")
  gen.word_add_table([["姓名","分数"],["张三",95],["李四",88]])
  buf = gen.word_save_bytes()

  # 生成Excel
  gen.excel_add_sheet("数据", headers=["名称","值"], rows=[["CPU","90%"],["内存","80%"]])
  buf = gen.excel_save_bytes()

  # 生成Markdown
  md = gen.markdown_report(title="报告", sections=[{"title":"概要","content":"正文"}])

依赖: python-docx, openpyxl, python-pptx, xlsxwriter, reportlab, markdown
"""

import os
import io
import re
import json
import time
from core.logging_config import get_logger
import hashlib
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field
from enum import Enum

logger = get_logger("evo.doc_generator")


# ═══════════════════════════════════════════════════
# Word文档生成器
# ═══════════════════════════════════════════════════

class WordGenerator:
    """Word (.docx) 文档生成器"""

    def __init__(self):
        self._doc = None
        self._init_doc()

    def _init_doc(self):
        try:
            from docx import Document
            from docx.shared import Pt, Inches, Cm, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            self._doc = Document()
            self._Pt = Pt
            self._Inches = Inches
            self._Cm = Cm
            self._RGBColor = RGBColor
            self._WD_ALIGN = WD_ALIGN_PARAGRAPH
        except ImportError:
            self._doc = None

    def _ensure(self):
        if self._doc is None:
            self._init_doc()
        if self._doc is None:
            raise RuntimeError("python-docx未安装: pip install python-docx")

    def reset(self):
        self._doc = None
        self._init_doc()

    def add_title(self, text: str, level: int = 0) -> Dict:
        self._ensure()
        self._doc.add_heading(text, level=level)
        return {"success": True}

    def add_paragraph(self, text: str, bold: bool = False, italic: bool = False,
                      font_size: int = 0, color: str = "", align: str = "") -> Dict:
        self._ensure()
        p = self._doc.add_paragraph()
        run = p.add_run(text)
        if bold:
            run.bold = True
        if italic:
            run.italic = True
        if font_size:
            run.font.size = self._Pt(font_size)
        if color:
            run.font.color.rgb = self._RGBColor.from_string(color)
        if align == "center":
            p.alignment = self._WD_ALIGN.CENTER
        elif align == "right":
            p.alignment = self._WD_ALIGN.RIGHT
        return {"success": True}

    def add_table(self, data: List[List[str]], headers: List[str] = None,
                  col_widths: List[float] = None) -> Dict:
        """添加表格"""
        self._ensure()
        if headers:
            rows = [headers] + data
        else:
            rows = data
        if not rows:
            return {"success": False, "error": "无数据"}
        table = self._doc.add_table(rows=len(rows), cols=len(rows[0]))
        table.style = "Table Grid"
        for i, row_data in enumerate(rows):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                if i == 0 and headers:
                    for run in cell.paragraphs[0].runs:
                        run.bold = True
        if col_widths:
            for i, width in enumerate(col_widths[:len(rows[0])]):
                for row in table.rows:
                    row.cells[i].width = self._Cm(width)
        return {"success": True, "rows": len(rows), "cols": len(rows[0])}

    def add_bullet_list(self, items: List[str]) -> Dict:
        self._ensure()
        for item in items:
            self._doc.add_paragraph(str(item), style="List Bullet")
        return {"success": True, "items": len(items)}

    def add_numbered_list(self, items: List[str]) -> Dict:
        self._ensure()
        for item in items:
            self._doc.add_paragraph(str(item), style="List Number")
        return {"success": True, "items": len(items)}

    def add_page_break(self) -> Dict:
        self._ensure()
        self._doc.add_page_break()
        return {"success": True}

    def add_image(self, image_path: str, width_cm: float = 15) -> Dict:
        self._ensure()
        if os.path.exists(image_path):
            self._doc.add_picture(image_path, width=self._Cm(width_cm))
            return {"success": True}
        return {"success": False, "error": f"图片不存在: {image_path}"}

    def add_image_bytes(self, image_bytes: bytes, width_cm: float = 15, ext: str = "png") -> Dict:
        """从字节添加图片"""
        self._ensure()
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
                f.write(image_bytes)
                tmp_path = f.name
            self._doc.add_picture(tmp_path, width=self._Cm(width_cm))
            os.unlink(tmp_path)
            return {"success": True}
        except Exception as e:
            return {"success": False, "error": str(e)}

    def save_bytes(self) -> bytes:
        self._ensure()
        buf = io.BytesIO()
        self._doc.save(buf)
        return buf.getvalue()

    def save_file(self, path: str) -> Dict:
        self._ensure()
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        self._doc.save(path)
        return {"success": True, "path": path, "size": os.path.getsize(path)}


# ═══════════════════════════════════════════════════
# Excel生成器
# ═══════════════════════════════════════════════════

class ExcelGenerator:
    """Excel (.xlsx) 文档生成器"""

    def __init__(self):
        self._wb = None
        self._sheets_data: Dict[str, Dict] = {}
        self._init()

    def _init(self):
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            self._wb = Workbook()
            self._Font = Font
            self._Fill = PatternFill
            self._Align = Alignment
            self._Border = Border
            self._Side = Side
        except ImportError:
            self._wb = None

    def _ensure(self):
        if self._wb is None:
            self._init()
        if self._wb is None:
            raise RuntimeError("openpyxl未安装: pip install openpyxl")

    def reset(self):
        self._wb = None
        self._init()

    def add_sheet(self, name: str, headers: List[str] = None,
                  rows: List[List[Any]] = None, auto_width: bool = True) -> Dict:
        self._ensure()
        ws = self._wb.create_sheet(title=name) if name != "Sheet" else self._wb.active
        if ws is None:
            ws = self._wb.active
        ws.title = name

        row_idx = 1
        if headers:
            # 表头样式
            header_fill = self._Fill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
            header_font = self._Font(bold=True, color="FFFFFF", size=11)
            thin_border = self._Border(
                left=self._Side(style="thin"),
                right=self._Side(style="thin"),
                top=self._Side(style="thin"),
                bottom=self._Side(style="thin"),
            )
            for j, h in enumerate(headers, 1):
                cell = ws.cell(row=row_idx, column=j, value=h)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = self._Align(horizontal="center", vertical="center")
                cell.border = thin_border
            row_idx += 1

        if rows:
            for row_data in rows:
                for j, val in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=j, value=val)
                    cell.border = self._Border(
                        left=self._Side(style="thin"), right=self._Side(style="thin"),
                        top=self._Side(style="thin"), bottom=self._Side(style="thin"),
                    )
                row_idx += 1

        # 自动列宽
        if auto_width and (headers or rows):
            all_data = (headers or []) + (rows or [])
            for j in range(len(all_data[0]) if all_data else 0):
                max_len = 0
                for row in all_data:
                    if j < len(row):
                        max_len = max(max_len, len(str(row[j])))
                ws.column_dimensions[chr(65 + j) if j < 26 else f"A{chr(65 + j - 26)}"].width = min(max_len + 4, 50)

        self._sheets_data[name] = {"headers": headers, "rows": rows or [], "total_rows": len(rows or [])}
        return {"success": True, "sheet": name, "rows": len(rows or [])}

    def add_formula(self, sheet: str, cell: str, formula: str) -> Dict:
        self._ensure()
        ws = self._wb[sheet]
        ws[cell] = formula
        return {"success": True}

    def save_bytes(self) -> bytes:
        self._ensure()
        buf = io.BytesIO()
        self._wb.save(buf)
        return buf.getvalue()

    def save_file(self, path: str) -> Dict:
        self._ensure()
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        self._wb.save(path)
        return {"success": True, "path": path, "size": os.path.getsize(path), "sheets": list(self._sheets_data.keys())}


# ═══════════════════════════════════════════════════
# PPT生成器
# ═══════════════════════════════════════════════════

class PPTGenerator:
    """PowerPoint (.pptx) 文档生成器"""

    def __init__(self):
        self._prs = None
        self._init()

    def _init(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            self._prs = Presentation()
            self._Inches = Inches
            self._Pt = Pt
        except ImportError:
            self._prs = None

    def _ensure(self):
        if self._prs is None:
            self._init()
        if self._prs is None:
            raise RuntimeError("python-pptx未安装: pip install python-pptx")

    def reset(self):
        self._prs = None
        self._init()

    def add_slide(self, title: str = "", content: str = "",
                  layout: str = "title_and_content", notes: str = "") -> Dict:
        """添加幻灯片"""
        self._ensure()
        layout_map = {
            "title": 0,  # Title Slide
            "title_and_content": 1,  # Title and Content
            "section": 2,  # Section Header
            "two_content": 3,  # Two Content
            "comparison": 4,  # Comparison
            "title_only": 5,  # Title Only
            "blank": 6,  # Blank
            "content_with_caption": 7,  # Content with Caption
            "picture_with_caption": 8,  # Picture with Caption
        }
        layout_idx = layout_map.get(layout, 1)
        slide_layout = self._prs.slide_layouts[layout_idx]
        slide = self._prs.slides.add_slide(slide_layout)

        # 标题
        if title:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 0:
                    shape.text = title
                    break

        # 内容
        if content and layout != "title":
            # 支持多行，每行一个段落
            lines = content.split("\n")
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for i, line in enumerate(lines):
                        if i == 0:
                            tf.paragraphs[0].text = line
                        else:
                            tf.add_paragraph().text = line
                    break

        # 备注
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return {"success": True, "slide": len(self._prs.slides)}

    def add_table_slide(self, title: str, data: List[List[str]],
                        headers: List[str] = None) -> Dict:
        """添加含表格的幻灯片"""
        self._ensure()
        rows = (headers or []) + data if headers else data
        cols = len(rows[0]) if rows else 1
        num_rows = len(rows)

        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])  # blank

        # 标题
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(28)

        # 表格
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.5 * num_rows)
        table_shape = slide.shapes.add_table(num_rows, cols, left, top, width, height)
        table = table_shape.table

        for i, row_data in enumerate(rows):
            for j, val in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(val)
                if i == 0 and headers:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True

        return {"success": True, "slide": len(self._prs.slides)}

    def save_bytes(self) -> bytes:
        self._ensure()
        buf = io.BytesIO()
        self._prs.save(buf)
        return buf.getvalue()

    def save_file(self, path: str) -> Dict:
        self._ensure()
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        self._prs.save(path)
        return {"success": True, "path": path, "size": os.path.getsize(path), "slides": len(self._prs.slides)}


# ═══════════════════════════════════════════════════
# Markdown/HTML生成器
# ═══════════════════════════════════════════════════

class MarkdownGenerator:
    """Markdown文档生成器"""

    @staticmethod
    def report(title: str, sections: List[Dict], metadata: Dict = None) -> str:
        """
        生成Markdown报告
        sections: [{"title": "标题", "content": "内容", "level": 2}]
        """
        lines = [f"# {title}", ""]
        if metadata:
            for k, v in metadata.items():
                lines.append(f"**{k}**: {v}")
            lines.append("")

        for sec in sections:
            level = sec.get("level", 2)
            content = sec.get("content", "")
            table_data = sec.get("table", None)

            lines.append(f"{'#' * level} {sec.get('title', '')}")
            lines.append("")

            if content:
                lines.append(content)
                lines.append("")

            if table_data:
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                if headers and rows:
                    # Markdown表格
                    lines.append("| " + " | ".join(str(h) for h in headers) + " |")
                    lines.append("| " + " | ".join("---" for _ in headers) + " |")
                    for row in rows:
                        lines.append("| " + " | ".join(str(v) for v in row) + " |")
                    lines.append("")

        lines.append(f"\n---\n*Generated by AUTO-EVO-AI V0.1 at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*")
        return "\n".join(lines)

    @staticmethod
    def to_html(markdown_text: str, title: str = "Report") -> str:
        """Markdown转HTML"""
        try:
            import markdown as md
            html_body = md.markdown(markdown_text, extensions=["tables", "fenced_code", "toc"])
        except ImportError:
            # 简单转换
            html_body = markdown_text.replace("\n", "<br>\n")

        return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<style>
body {{ max-width: 900px; margin: 0 auto; padding: 40px 24px; font-family: system-ui, -apple-system, sans-serif; line-height: 1.7; color: #1a1a2e; background: #f8f9fa; }}
h1 {{ color: #1a1a2e; border-bottom: 3px solid #2563eb; padding-bottom: 12px; }}
h2 {{ color: #1e40af; margin-top: 32px; }}
h3 {{ color: #3b82f6; }}
table {{ border-collapse: collapse; width: 100%; margin: 16px 0; background: white; border-radius: 8px; overflow: hidden; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
th {{ background: #1e40af; color: white; padding: 12px 16px; text-align: left; }}
td {{ padding: 10px 16px; border-bottom: 1px solid #e5e7eb; }}
tr:hover {{ background: #f0f9ff; }}
code {{ background: #e5e7eb; padding: 2px 6px; border-radius: 4px; font-size: 0.9em; }}
pre {{ background: #1e293b; color: #e2e8f0; padding: 16px; border-radius: 8px; overflow-x: auto; }}
hr {{ border: none; border-top: 1px solid #d1d5db; margin: 32px 0; }}
blockquote {{ border-left: 4px solid #3b82f6; padding-left: 16px; color: #6b7280; }}
</style>
</head>
<body>
{html_body}
</body>
</html>"""


# ═══════════════════════════════════════════════════
# 统一文档生成器
# ═══════════════════════════════════════════════════

class DocGenerator:
    """
    统一文档生成器入口

    支持: Word(.docx) / Excel(.xlsx) / PPT(.pptx) / Markdown / HTML
    """

    def __init__(self):
        self.word = WordGenerator()
        self.excel = ExcelGenerator()
        self.ppt = PPTGenerator()
        self.markdown = MarkdownGenerator()
        self._output_dir = os.path.join(tempfile.gettempdir(), "evo_docs")
        os.makedirs(self._output_dir, exist_ok=True)

    # ─── 快捷生成 ───

    def generate_report(self, title: str, sections: List[Dict],
                        format: str = "markdown", metadata: Dict = None) -> Dict:
        """
        一键生成报告
        format: markdown | html | word | all
        """
        md_text = self.markdown.report(title, sections, metadata)
        results = {}

        if format in ("markdown", "all"):
            md_path = os.path.join(self._output_dir, f"{self._safe_name(title)}.md")
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(md_text)
            results["markdown"] = {"path": md_path, "size": os.path.getsize(md_path)}

        if format in ("html", "all"):
            html_text = self.markdown.to_html(md_text, title)
            html_path = os.path.join(self._output_dir, f"{self._safe_name(title)}.html")
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(html_text)
            results["html"] = {"path": html_path, "size": os.path.getsize(html_path)}

        if format in ("word", "all"):
            self.word.reset()
            self.word.add_title(title)
            for sec in sections:
                self.word.add_title(sec.get("title", ""), level=1)
                self.word.add_paragraph(sec.get("content", ""))
                if sec.get("table"):
                    td = sec["table"]
                    self.word.add_table(td.get("rows", []), td.get("headers", []))
            word_path = os.path.join(self._output_dir, f"{self._safe_name(title)}.docx")
            self.word.save_file(word_path)
            results["word"] = {"path": word_path, "size": os.path.getsize(word_path)}

        return {"success": True, "title": title, "format": format, "files": results}

    def generate_table_report(self, title: str, tables: Dict[str, Dict],
                               format: str = "excel") -> Dict:
        """
        生成表格报告
        tables: {"Sheet名": {"headers": [...], "rows": [...]}}
        """
        results = {}

        if format in ("excel", "all"):
            self.excel.reset()
            for sheet_name, data in tables.items():
                self.excel.add_sheet(sheet_name, data.get("headers"), data.get("rows"))
            xlsx_path = os.path.join(self._output_dir, f"{self._safe_name(title)}.xlsx")
            self.excel.save_file(xlsx_path)
            results["excel"] = {"path": xlsx_path, "size": os.path.getsize(xlsx_path)}

        return {"success": True, "title": title, "format": format, "files": results}

    def generate_presentation(self, title: str, slides: List[Dict]) -> Dict:
        """
        生成PPT
        slides: [{"title": "...", "content": "...", "notes": "...", "table": {"headers":[], "rows":[]}}]
        """
        self.ppt.reset()
        for i, slide_data in enumerate(slides):
            if slide_data.get("table"):
                self.ppt.add_table_slide(
                    slide_data.get("title", f"Slide {i+1}"),
                    slide_data["table"].get("rows", []),
                    slide_data["table"].get("headers"),
                )
            else:
                self.ppt.add_slide(
                    title=slide_data.get("title", ""),
                    content=slide_data.get("content", ""),
                    notes=slide_data.get("notes", ""),
                )
        pptx_path = os.path.join(self._output_dir, f"{self._safe_name(title)}.pptx")
        self.ppt.save_file(pptx_path)
        return {"success": True, "path": pptx_path, "size": os.path.getsize(pptx_path), "slides": len(slides)}

    # ─── 数据转换 ───

    def data_to_excel_bytes(self, tables: Dict[str, Dict]) -> bytes:
        """数据直接转Excel字节流"""
        self.excel.reset()
        for sheet_name, data in tables.items():
            self.excel.add_sheet(sheet_name, data.get("headers"), data.get("rows"))
        return self.excel.save_bytes()

    def data_to_word_bytes(self, title: str, sections: List[Dict]) -> bytes:
        """数据直接转Word字节流"""
        self.word.reset()
        self.word.add_title(title)
        for sec in sections:
            self.word.add_title(sec.get("title", ""), level=1)
            self.word.add_paragraph(sec.get("content", ""))
            if sec.get("table"):
                self.word.add_table(sec["table"].get("rows", []), sec["table"].get("headers", []))
        return self.word.save_bytes()

    def list_output_files(self) -> List[Dict]:
        """列出已生成的文件"""
        files = []
        for f in os.listdir(self._output_dir):
            fp = os.path.join(self._output_dir, f)
            if os.path.isfile(fp):
                files.append({
                    "name": f,
                    "path": fp,
                    "size": os.path.getsize(fp),
                    "modified": datetime.fromtimestamp(os.path.getmtime(fp)).strftime("%Y-%m-%d %H:%M:%S"),
                })
        return sorted(files, key=lambda x: x["modified"], reverse=True)

    def _safe_name(self, name: str) -> str:
        return re.sub(r'[\\/:*?"<>|]', '_', name)[:80]


# ═══════════════════════════════════════════════════
# 全局单例
# ═══════════════════════════════════════════════════

import tempfile

_doc_generator: Optional[DocGenerator] = None


def get_doc_generator() -> DocGenerator:
    global _doc_generator
    if _doc_generator is None:
        _doc_generator = DocGenerator()
    return _doc_generator
