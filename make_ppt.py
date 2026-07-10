"""AUTO-EVO-AI 建筑玻璃膜介绍 PPT — 5页"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 调色板
BG = RGBColor(0x1A,0x1A,0x2E)
W = RGBColor(0xFF,0xFF,0xFF)
A = RGBColor(0x43,0x61,0xEE)
G = RGBColor(0x72,0x09,0xB7)
M = RGBColor(0x88,0x92,0xB0)
CARD = RGBColor(0x16,0x21,0x3E)
BORDER = RGBColor(0x2D,0x35,0x61)

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_text(slide, left, top, w, h, text, size=18, bold=False, color=W, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb

def add_bullet(slide, left, top, w, h, items, size=16, color=W):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)

# ═══════════ 第1页：封面 ═══════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_shape(sl, Inches(1), Inches(1.8), Inches(11.3), Inches(3.8), CARD)
add_text(sl, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
         "建筑玻璃膜介绍", 48, True, W, PP_ALIGN.CENTER)
add_text(sl, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
         "提升建筑美学与功能性的创新材料", 22, False, M, PP_ALIGN.CENTER)
add_text(sl, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
         "AUTO-EVO-AI  |  2026", 16, False, M, PP_ALIGN.CENTER)

# ═══════════ 第2页：什么是建筑玻璃膜 ═══════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "什么是建筑玻璃膜？", 36, True, A)
add_shape(sl, Inches(0.8), Inches(1.3), Inches(0.6), Inches(0.06), A)

add_shape(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), CARD)
add_bullet(sl, Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2), [
    "建筑玻璃膜是一种贴在玻璃表面的高性能薄膜材料",
    "由聚酯（PET）基材与多层光学涂层复合而成",
    "厚度仅0.05-0.3mm，几乎不增加玻璃重量",
    "广泛应用于商业建筑、住宅、公共设施等",
    "兼具安全防护、节能隔热、隐私保护等功能",
    "施工便捷，成本远低于更换玻璃",
], 17, W)

add_shape(sl, Inches(7), Inches(1.8), Inches(5.5), Inches(4.8), CARD)
add_bullet(sl, Inches(7.4), Inches(2.1), Inches(4.8), Inches(4.2), [
    "玻璃膜起源于20世纪60年代的航天技术",
    "最初用于卫星和航天器的热控系统",
    "1970年代开始应用于建筑领域",
    "如今已成为绿色建筑标准配置之一",
    "全球市场规模超200亿美元且持续增长",
    "中国建筑玻璃膜渗透率仅5%，潜力巨大",
], 17, W)

# ═══════════ 第3页：玻璃膜的类型 ═══════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "玻璃膜的类型", 36, True, A)
add_shape(sl, Inches(0.8), Inches(1.3), Inches(0.6), Inches(0.06), A)

cards = [
    ("🛡️ 安全膜", [
        "防爆防弹，抗冲击性强",
        "玻璃破碎后粘合成片",
        "防止碎片飞溅伤人",
        "抗风压、抗震能力提升",
        "符合国家安全玻璃标准",
        "厚度：0.1-0.3mm",
    ]),
    ("☀️ 隔热膜", [
        "阻隔90%以上紫外线",
        "红外线阻隔率高达85%",
        "夏季降温3-8°C",
        "冬季保温减少热流失",
        "节能率达30%-50%",
        "降低空调电费支出",
    ]),
    ("🎨 装饰膜", [
        "磨砂、压花等多种纹理",
        "仿古铜、木纹等效果",
        "个性化定制图案色彩",
        "保护隐私的同时透光",
        "办公室隔断首选方案",
        "可随时更换不留胶",
    ]),
]

for i, (title, items) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    add_shape(sl, x, Inches(2), Inches(3.8), Inches(5), CARD)
    add_shape(sl, x, Inches(2), Inches(3.8), Inches(0.7), A)
    add_text(sl, x, Inches(2.1), Inches(3.8), Inches(0.6),
             title, 22, True, W, PP_ALIGN.CENTER)
    add_bullet(sl, Inches(0.8 + i*4.1 + 0.3), Inches(3), Inches(3.2), Inches(3.8),
               items, 15, W)

# ═══════════ 第4页：玻璃膜的优势 ═══════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "玻璃膜的四大核心优势", 36, True, A)
add_shape(sl, Inches(0.8), Inches(1.3), Inches(0.6), Inches(0.06), A)

advantages = [
    ("⚡ 节能环保", "阻隔90%紫外线+85%红外线\n夏季降温3-8°C，降低空调能耗30-50%\n减少碳排放，助力绿色建筑认证", RGBColor(0x43,0x61,0xEE)),
    ("🛡️ 安全保障", "防爆防弹等级可选\n玻璃破碎后粘合成片，不飞溅\n抗风压、抗震、防盗性能提升", RGBColor(0x72,0x09,0xB7)),
    ("🎭 隐私美观", "从外到内遮蔽视线\n从内到外清晰可见\n多种纹理色彩可选\n提升建筑整体美感", RGBColor(0x02,0x80,0x90)),
    ("💰 经济实用", "成本仅为更换玻璃的10-30%\n施工快速，不影响正常办公\n使用寿命8-15年\n投资回报周期仅2-3年", RGBColor(0xE6,0x5C,0x00)),
]

for i, (title, desc, color) in enumerate(advantages):
    x = Inches(0.8 + i * 3.1)
    add_shape(sl, x, Inches(2), Inches(2.8), Inches(5), CARD)
    add_shape(sl, x, Inches(2), Inches(2.8), Inches(0.6), color)
    add_text(sl, x, Inches(2.1), Inches(2.8), Inches(0.5),
             title, 18, True, W, PP_ALIGN.CENTER)
    add_text(sl, Inches(0.8 + i*3.1 + 0.2), Inches(2.9),
             Inches(2.4), Inches(3.8),
             desc, 14, False, M)

# ═══════════ 第5页：案例与应用 ═══════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text(sl, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "案例分析与未来趋势", 36, True, A)
add_shape(sl, Inches(0.8), Inches(1.3), Inches(0.6), Inches(0.06), A)

add_shape(sl, Inches(0.8), Inches(2), Inches(5.5), Inches(5), CARD)
add_text(sl, Inches(1.2), Inches(2.2), Inches(4.8), Inches(0.5),
         "📌 典型应用案例", 22, True, A)
add_bullet(sl, Inches(1.2), Inches(2.8), Inches(4.8), Inches(4), [
    "🏢 上海中心大厦 — 全楼隔热膜降低能耗35%",
    "🏨 北京国贸大酒店 — 安全膜达到防弹标准",
    "🏪 万达广场 — 装饰膜实现统一品牌视觉",
    "🏥 协和医院 — 紫外线阻隔保护药品存储",
    "🏫 清华大学图书馆 — 隔热膜保护古籍藏书",
    "🏠 万科翡翠 — 住宅项目隐私膜标配交付",
], 16, W)

add_shape(sl, Inches(7), Inches(2), Inches(5.5), Inches(5), CARD)
add_text(sl, Inches(7.4), Inches(2.2), Inches(4.8), Inches(0.5),
         "📈 未来发展趋势", 22, True, A)
add_bullet(sl, Inches(7.4), Inches(2.8), Inches(4.8), Inches(4), [
    "智能调光玻璃膜 — 电压控制透明度",
    "太阳能发电膜 — 让窗户变成发电站",
    "5G通讯兼容膜 — 不屏蔽信号的新型膜",
    "自清洁纳米涂层 — 减少人工维护成本",
    "建筑一体化BIPV — 幕墙本身就是能源",
    "AI驱动的定制化 — 算法设计最优膜层",
], 16, W)

# 保存
out = __file__.replace(".py", ".pptx")
prs.save(out)
logger.info(f"✅ PPT 已生成: {out}"))
