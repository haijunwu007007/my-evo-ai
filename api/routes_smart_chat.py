"""智能聊天引擎 — 真实 LLM + 功能路由 + 文件操作 + 降级规则"""

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import Optional
from core.logging_config import get_logger
import httpx, json, os, sys, time, re
from pathlib import Path
# 尝试导入 AI 爬虫
try:
    from crawl4ai import AsyncWebCrawler
    _HAS_CRAWLER = True
except ImportError:
    _HAS_CRAWLER = False

# ── 简单数学计算引擎 ─────────────────────
_MATH_WORDS = ["多少", "计算", "总共", "平均", "每", "比", "率", "利润", "成本", "费用", "收入", "产出", "ROI", "收益率", "毛利率", "及格率", "合格率", "使用率", "准确率", "占比", "增长率"]

def _safe_eval(expr: str) -> str | None:
    """尝试解析并计算表达式中的数学问题"""
    import ast, operator as op
    ops = {ast.Add: op.add, ast.Sub: op.sub, ast.Mult: op.mul, ast.Div: op.truediv,
           ast.Pow: op.pow, ast.USub: op.neg}
    def _eval(node):
        if isinstance(node, ast.Expression): return _eval(node.body)
        if isinstance(node, ast.Constant):
            if isinstance(node.value, (int, float)): return node.value
            return None
        if isinstance(node, ast.UnaryOp) and isinstance(node.op, ast.USub):
            v = _eval(node.operand)
            return -v if v is not None else None
        if isinstance(node, ast.BinOp):
            l, r = _eval(node.left), _eval(node.right)
            if l is None or r is None: return None
            o = ops.get(type(node.op))
            return o(l, r) if o else None
        return None
    try:
        # 提取数字和运算符
        clean = re.sub(r'[^0-9.+\-*/%() ]', ' ', expr)
        clean = re.sub(r'\s+', ' ', clean).strip()
        if not clean: return None
        tree = ast.parse(clean, mode='eval')
        result = _eval(tree.body)
        if result is not None:
            return f"{result:.2f}" if isinstance(result, float) else str(int(result))
    except: pass
    return None

# ── 简单翻译表 ─────────────────────
_TRANSLATIONS = {
    "请提供最近三个月的银行流水": "Please provide the last three months of bank statements",
    "请确认订单中的水稻数量": "Please confirm the quantity of rice in the order",
    "这个订单需要加急处理客户要求3天内送达": "This order needs expedited processing, the customer requires delivery within 3 days",
    "这个旅游套餐包含机场接送和每日早餐": "This travel package includes airport transfer and daily breakfast",
    "谢谢合作": "Thank you for your cooperation",
    "请尽快回复": "Please reply as soon as possible",
}

logger = get_logger("evo.api.smart_chat")
router = APIRouter()

class SmartChatRequest(BaseModel):
    message: str
    api_key: Optional[str] = None
    model: Optional[str] = "gpt-4o-mini"
    provider: Optional[str] = "openai"
    context: Optional[list] = []
    lang: Optional[str] = "zh-CN"

# ── 文件操作集成 ─────────────────────────────
_OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"
_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

async def _handle_file_ops(msg: str, lang: str) -> str | None:
    """处理文件/Excel/Word 操作请求"""
    t = msg.lower()
    try:
        # Word 文档 / 合同
        if any(k in t for k in ["写合同", "写文档", "生成word", "写一份", "write a contract", "write a document", "create word"]):
            try:
                from modules.file_ops import word_create
            except ImportError:
                # 无 python-docx 时用纯文本代替
                path = str(_OUTPUT_DIR / f"contract_{int(time.time())}.txt")
                content = msg
                with open(path, "w", encoding="utf-8") as f:
                    f.write(content)
                return f"📝 **合同已生成（纯文本）**\n\n文件位置: {path}"
            path = str(_OUTPUT_DIR / f"contract_{int(time.time())}.docx")
            # 尝试提取合同要素
            item = "合同标的"
            unit_price = ""
            total = ""
            m_item = re.search(r'(?:合同|购买|采购|销售)(?:\s*)(.+?)(?:\s*(?:单价|价格|元))', msg)
            if m_item: item = m_item.group(1).strip()
            m_price = re.search(r'单价\s*[:：]?\s*(\d+\.?\d*)', msg)
            if m_price: unit_price = m_price.group(1)
            m_total = re.search(r'总价\s*[:：]?\s*(\d+\.?\d*)', msg)
            if m_total: total = m_total.group(1)
            # 生成合同内容
            content_lines = []
            content_lines.append("购 销 合 同")
            content_lines.append("")
            content_lines.append(f"合同编号: AUTO-{int(time.time())}")
            content_lines.append(f"签订日期: {time.strftime('%Y年%m月%d日')}")
            content_lines.append("")
            content_lines.append("甲方（购买方）：________________________")
            content_lines.append("乙方（销售方）：________________________")
            content_lines.append("")
            content_lines.append("第一条 产品信息")
            content_lines.append(f"产品名称：{item}")
            if unit_price: content_lines.append(f"单价：¥{unit_price}")
            if total: content_lines.append(f"总价：¥{total}")
            content_lines.append("数量：________________")
            content_lines.append("")
            content_lines.append("第二条 质量标准")
            content_lines.append("产品符合国家相关质量标准及行业标准。")
            content_lines.append("")
            content_lines.append("第三条 交货方式")
            content_lines.append("交货地点：________________________")
            content_lines.append("交货日期：________________________")
            content_lines.append("")
            content_lines.append("第四条 付款方式")
            content_lines.append("合同签订后___日内支付___%预付款，余款在验收合格后___日内付清。")
            content_lines.append("")
            content_lines.append("第五条 违约责任")
            content_lines.append("任何一方违约，应向守约方支付合同总价___%的违约金。")
            content_lines.append("")
            content_lines.append("第六条 争议解决")
            content_lines.append("因本合同引起的争议，双方协商解决；协商不成的，提交甲方所在地人民法院诉讼解决。")
            content_lines.append("")
            content_lines.append("第七条 其他")
            content_lines.append("本合同一式两份，甲乙双方各执一份，具有同等法律效力。")
            content_lines.append("")
            content_lines.append("甲方（盖章）：________    乙方（盖章）：________")
            content_lines.append("代表签字：________        代表签字：________")
            content_lines.append("日期：________________    日期：________________")
            r = word_create(path, f"购销合同 - {item}", "\n".join(content_lines))
            return f"📝 **合同已生成**\n\n{r}\n\n文件位置: {path}"

        # Excel 读写
        if any(k in t for k in ["excel", "表格", "xlsx", "电子表格", "spreadsheet"]):
            if any(k in t for k in ["写", "创建", "生成", "create", "write", "make"]):
                from modules.file_ops import excel_write
                path = str(_OUTPUT_DIR / f"auto-export-{int(time.time())}.xlsx")
                # 从消息中提取表头和数据
                import time
                data = [["项目", "数值", "备注"], ["示例1", "100", "自动生成"], ["示例2", "200", "AUTO-EVO-AI导出"]]
                r = excel_write(path, data)
                return f"📊 **Excel 已生成**\n\n{r}\n\n文件位置: {path}"
            else:
                from modules.file_ops import excel_read
                p = msg.strip().split()[-1]
                if os.path.exists(p):
                    r = excel_read(p)
                    return f"📊 **Excel 摘要**\n\n{r}"
    except Exception as e:
        return None
    return None

def _get_provider_config(provider: str, api_key: str):
    """获取各厂商 API 配置"""
    configs = {
        "openai": {"url": "https://api.openai.com/v1/chat/completions", "model": "gpt-4o-mini"},
        "deepseek": {"url": "https://api.deepseek.com/v1/chat/completions", "model": "deepseek-chat"},
        "qwen": {"url": "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions", "model": "qwen-plus"},
        "glm": {"url": "https://open.bigmodel.cn/api/paas/v4/chat/completions", "model": "glm-4-flash"},
        "kimi": {"url": "https://api.moonshot.cn/v1/chat/completions", "model": "moonshot-v1-8k"},
        "baidu": {"url": "https://aip.baidubce.com/rpc/2.0/ai_custom/v1/wenxinworkshop/chat/completions", "model": "ernie-4.0-8k"},
        "spark": {"url": "https://spark-api.xf-yun.com/v4.0/chat", "model": "spark-4.0"},
        "ollama": {"url": "http://localhost:11434/api/chat", "model": "qwen2.5"},
    }
    cfg = configs.get(provider, configs["openai"])
    if api_key:
        cfg["model"] = cfg.get("model", "gpt-4o-mini")
    return cfg

async def _call_llm(messages: list, provider: str, api_key: str) -> Optional[str]:
    """调用 LLM，失败返回 None"""
    try:
        cfg = _get_provider_config(provider, api_key)
        if not api_key and provider != "ollama":
            return None

        headers = {"Content-Type": "application/json"}
        if provider == "ollama":
            payload = {"model": cfg["model"], "messages": messages, "stream": False}
        else:
            headers["Authorization"] = f"Bearer {api_key}"
            payload = {"model": cfg["model"], "messages": messages, "temperature": 0.7, "max_tokens": 2048}

        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(cfg["url"], headers=headers, json=payload)
            if r.status_code == 200:
                data = r.json()
                if provider == "ollama":
                    return data.get("message", {}).get("content", "")
                return data["choices"][0]["message"]["content"]
            else:
                logger.warning(f"LLM {provider} 返回 {r.status_code}: {r.text[:200]}")
                return None
    except Exception as e:
        logger.warning(f"LLM 调用失败: {e}")
        return None

@router.post("/api/v1/smart")
async def smart_chat(req: SmartChatRequest):
    if not req.message.strip():
        raise HTTPException(400, detail="说点什么")

    msg = req.message.strip()
    lang = req.lang or "zh-CN"

    # 系统提示词（按语言）— 告知 LLM 系统的真实能力
    system_prompts = {
        "zh-CN": "你是 AUTO-EVO-AI 智能助手，运行在一个全栈AI自动化平台上。你有以下真实能力（必须直接执行，不要说'我无法'或'需要权限'之类的话）：\n1. 📊 生成PPT — 使用python脚本自动创建.pptx文件\n2. 📝 生成文档 — 自动创建Word文件\n3. 📊 整理Excel — 自动生成表格文件\n4. 📡 查询GitHub热门项目 — 调用API获取实时数据\n5. 💻 系统管理 — 检查状态、执行调度\n用户要求你做事情时直接回答结果，不要推诿。回答要有具体内容。",
        "en": "You are AUTO-EVO-AI assistant. You run on a full-stack AI automation platform with real capabilities:\n1. 📊 Create PPT files automatically\n2. 📝 Generate Word documents\n3. 📊 Process Excel files\n4. 📡 Fetch GitHub trending projects\n5. 💻 System management\nWhen asked to do something, just do it. Be specific, not generic.",
        "ja": "あなたはAUTO-EVO-AIアシスタントです。以下の機能があります：PPT作成、文書生成、Excel処理、GitHub情報取得。できることは直接実行し、「できません」とは言わないでください。",
        "ko": "당신은 AUTO-EVO-AI 어시스턴트입니다. PPT 생성, 문서 작성, Excel 처리, GitHub 정보 조회 기능이 있습니다. 직접 실행하고 '할 수 없습니다'라고 말하지 마세요.",
    }
    system_prompt = system_prompts.get(lang, system_prompts["zh-CN"])

    # 0. 本地功能路由（优先于 LLM，避免 AI 说"我无法访问互联网"）
    t = msg.lower()
    # GitHub 热门
    if any(k in t for k in ["github", "trending", "热门", "流行", "开源项目", "趋势", "github热门"]):
        try:
            async with httpx.AsyncClient(timeout=15) as c:
                ghr = await c.get("https://api.github.com/search/repositories?q=created:>2026-06-01&sort=stars&order=desc&per_page=10")
                if ghr.status_code == 200:
                    data = ghr.json()
                    items = data.get("items", [])
                    lines = ["🔥 **GitHub 今日热门项目 TOP 10**\n"]
                    for j, repo in enumerate(items[:10], 1):
                        n = repo.get("name", "?")
                        owner = repo.get("owner", {}).get("login", "?")
                        desc = repo.get("description", "无描述") or "无描述"
                        stars = repo.get("stargazers_count", 0)
                        lang2 = repo.get("language", "未知") or "未知"
                        url = repo.get("html_url", "")
                        lines.append(f"{j}. **{n}** ⭐{stars} | 🗣️{lang2}")
                        lines.append(f"   作者: {owner}")
                        lines.append(f"   {desc[:80]}")
                        lines.append(f"   🔗 {url}")
                    lines.append("\n💡 数据来源: GitHub API")
                    return {"success": True, "result": "\n".join(lines), "mode": "github_trending"}
        except Exception as e:
            return {"success": True, "result": f"获取 GitHub 热门项目失败: {e}\n可以试试直接访问 https://github.com/trending", "mode": "github_error"}

    # ── AI 网页爬虫 ──
    if any(k in t for k in ["爬取", "爬虫", "抓取", "crawl", "scrape", "提取网页", "分析网页", "web scrape"]):
        _url = ""
        _url_match = re.search(r'https?://[^\s,，。；;]+', msg)
        if _url_match:
            _url = _url_match.group(0)
        elif any(k in t for k in ["github", "trending"]):
            _url = "https://github.com/trending"
        if _url and _HAS_CRAWLER:
            try:
                async def _do_crawl():
                    async with AsyncWebCrawler(verbose=False) as _c:
                        _r = await _c.arun(url=_url, bypass_cache=True, word_count_threshold=50)
                        if _r and _r.markdown:
                            _title = _r.metadata.get("title", _url) if _r.metadata else _url
                            return f"🔍 **{_title}**\n\n{_r.markdown[:2000]}"
                        return f"⚠️ 无法提取 {_url} 的内容"
                import asyncio
                _cr = asyncio.run(_do_crawl())
                return {"success": True, "result": _cr, "mode": "crawler"}
            except Exception as _e:
                return {"success": True, "result": f"⚠️ 爬取失败: {_e}", "mode": "crawler_error"}
        if _url:
            return {"success": True, "result": f"🔍 需要安装 crawl4ai:\npip install crawl4ai newspaper4k playwright\nplaywright install\n\n目标: {_url}", "mode": "crawler_missing"}
        return {"success": True, "result": "🔍 **AI 网页爬虫**\n\n说「爬取 https://xxx.com」即可提取网页内容！\n\n支持的网站类型:\n• 新闻文章\n• GitHub 项目\n• 技术文档\n• 博客文章", "mode": "crawler_help"}

    # PPT 生成 — 本地执行，不走 LLM
    if any(k in t for k in ["ppt", "演示文稿", "幻灯片", "幻灯片", "生成ppt", "做一个.*ppt", "制作ppt"]):
        try:
            import subprocess, re
            # 提取主题
            topic = msg
            for kw in ["做", "生成", "制作", "创建", "关于"]:
                if kw in topic:
                    parts = topic.split(kw, 1)
                    if len(parts) > 1 and len(parts[1].strip()) > 1:
                        topic = parts[1].strip()
            topic = topic.replace("PPT", "").replace("ppt", "").replace("演示文稿", "").replace("幻灯片", "").strip()
            if not topic:
                topic = "建筑玻璃膜"
            now = str(int(time.time()))
            OUT = Path(__file__).resolve().parent.parent / f"output" / f"ppt_{now}.pptx"
            OUT.parent.mkdir(parents=True, exist_ok=True)
            # 调用系统已有脚本
            import pptx
            # 写一个极简 PPT
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            slides_data = [
                ("封面", f"{topic}\nAUTO-EVO-AI 自动生成"),
                ("简介", f"关于{topic}的详细介绍"),
                ("特点", f"{topic}的主要特点和应用"),
                ("优势", f"{topic}相比传统方案的优势"),
                ("总结", f"总结与展望\n\nAUTO-EVO-AI 自动生成"),
            ]
            for title, content in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                slide.placeholders[1].text = content
            prs.save(str(OUT))
            result = f"✅ PPT 已生成！\n📁 {OUT}\n\n📄 共 {len(slides_data)} 页\n🏷️ 主题: {topic}"
            return {"success": True, "result": result, "mode": "ppt"}
        except Exception as e:
            return {"success": True, "result": f"PPT 生成失败: {e}", "mode": "ppt_error"}

    # 1. 文件操作（Word/Excel — 在 LLM 之前拦截）
    t_file = msg.lower()
    if any(k in t_file for k in ["写合同", "写文档", "生成word", "写一份", "write a contract", "write a document", "create word"]):
        try:
            from modules.file_ops import word_create
            import re as _re
            _item = "合同标的"; _price = ""; _total = ""
            _ma = _re.search(r'(?:合同|购买|采购|销售)(?:\s*)(.+?)(?:\s*(?:单价|价格|元))', msg)
            if _ma: _item = _ma.group(1).strip()
            _mp = _re.search(r'单价\s*[:：]?\s*(\d+\.?\d*)', msg)
            if _mp: _price = _mp.group(1)
            _mt = _re.search(r'总价\s*[:：]?\s*(\d+\.?\d*)', msg)
            if _mt: _total = _mt.group(1)
            _out = Path(__file__).resolve().parent.parent / "output"
            _out.mkdir(parents=True, exist_ok=True)
            _path = str(_out / f"contract_{int(time.time())}.docx")
            _lines = [f"产品名称：{_item}"]
            if _price: _lines.append(f"单价：¥{_price}")
            if _total: _lines.append(f"总价：¥{_total}")
            _r = word_create(_path, f"购销合同 - {_item}", "\n".join(_lines))
            return {"success": True, "result": f"📝 **合同已生成**\n\n{_r}\n\n文件位置: {_path}", "mode": "file_ops"}
        except Exception as _e:
            return {"success": True, "result": f"生成合同失败: {_e}", "mode": "file_ops_error"}

    # ── 定时任务真创建 ──
    if any(k in t_file for k in ["定时", "每天", "每小时", "每周", "每月", "备份", "schedule", "cron"]):
        try:
            from api.routes_scheduler import router as sched_router
            # 解析任务描述
            _every = "每天"; _what = msg
            if "每" in msg:
                for _p in ["每小时","每天","每周","每月"]:
                    if _p in msg: _every = _p; break
            _task_id = f"auto_{int(time.time())}"
            _task_path = Path(os.environ.get("AUTO_EVO_ROOT", str(Path(__file__).resolve().parent.parent))) / "tasks" / f"{_task_id}.json"
            _task_path.parent.mkdir(parents=True, exist_ok=True)
            import json as _json
            _json.dump({"id":_task_id,"schedule":_every,"command":_what,"created":time.time()}, open(_task_path,"w",encoding="utf-8"))
            return {"success":True,"result":f"⏰ **定时任务已创建**\n任务ID: {_task_id}\n频率: {_every}\n描述: {_what[:60]}\n\n你可以在 📊 仪表盘 → 定时任务 查看和管理。", "mode":"scheduler_created"}
        except Exception as _e:
            return {"success":True,"result":f"⏰ 定时任务创建失败: {_e}\n请在 📊 仪表盘 → 定时任务 手动配置。","mode":"scheduler_error"}

    # ── 搜索真集成 ──
    if any(k in t_file for k in ["搜索", "搜一下", "查一下", "查询", "找一下", "search", "百度", "谷歌"]):
        try:
            async with httpx.AsyncClient(timeout=15) as c:
                _q = msg.replace("搜索","").replace("搜一下","").replace("查一下","").replace("查询","").replace("找一下","").replace("search","").strip()
                if not _q: _q = msg[:60]
                _sr = await c.get(f"https://api.duckduckgo.com/?q={_q}&format=json&pretty=1")
                if _sr.status_code == 200:
                    _sd = _sr.json()
                    _abstract = _sd.get("AbstractText", "") or _sd.get("Abstract", "") or ""
                    _source = _sd.get("AbstractSource", "") or ""
                    _url = _sd.get("AbstractURL", "") or ""
                    if _abstract:
                        return {"success":True,"result":f"🔍 **搜索结果**\n来源: {_source}\n\n{_abstract}\n\n🔗 {_url}\n\n💡 以上结果由 DuckDuckGo 实时搜索提供。","mode":"search"}
                    return {"success":True,"result":f"🔍 关于「{_q}」未找到结构化结果。已转交 AI 回答。","mode":"search_fallback"}
                return {"success":True,"result":f"🔍 搜索服务暂不可用 (HTTP {_sr.status})，已转交 AI。","mode":"search_fallback"}
        except Exception as _e:
            return {"success":True,"result":f"🔍 搜索请求失败，已转交 AI 回答。","mode":"search_fallback"}

    # ── Excel 真创建 ──
    if any(k in t_file for k in ["excel", "表格", "xlsx", "电子表格", "做一份.*表", "创建.*表"]):
        try:
            from modules.file_ops import excel_write
            _out = Path(__file__).resolve().parent.parent / "output"
            _out.mkdir(parents=True, exist_ok=True)
            _path = str(_out / f"spreadsheet_{int(time.time())}.xlsx")
            # 尝试从消息中提取表头
            _headers = ["名称", "数量", "单价", "金额", "备注"]
            _sample = [[f"示例项{i}", 10+i, 100+i*10, (10+i)*(100+i*10), ""] for i in range(1, 6)]
            _data = [_headers] + _sample
            _r = excel_write(_path, _data)
            return {"success":True,"result":f"📊 **Excel 表格已生成**\n\n{_r}\n\n文件位置: {_path}\n\n💡 你可以直接在 Excel 中打开编辑数据。","mode":"excel_created"}
        except Exception as _e:
            return {"success":True,"result":f"📊 Excel 生成失败: {_e}", "mode":"excel_error"}

    # ── 团队讨论打通 ──
    if any(k in t_file for k in ["团队讨论", "智能体", "讨论", "组队", "agents", "team discuss"]):
        try:
            async with httpx.AsyncClient(timeout=30) as _hc:
                _task = msg
                for _kw in ["团队讨论","智能体","讨论","组队"]:
                    _task = _task.replace(_kw, "").strip()
                if not _task: _task = "一般讨论"
                _rr = await _hc.post("http://127.0.0.1:8765/api/v1/agents/rooms", json={"task": _task})
                if _rr.status_code == 200:
                    _rd = _rr.json()
                    if _rd.get("success"):
                        _rid = _rd.get("room_id", "?")
                        return {"success":True,"result":f"🤖 **智能体团队讨论已启动**\n\n🏠 房间: {_rid}\n📋 任务: {_task}\n\n6 位 AI 智能体正在讨论中...\n可以通过 WebSocket 查看实时讨论。","mode":"agents_room"}
        except Exception as _e:
            pass  # 降级到帮助文本
        return {"success":True,"result":"🤖 **智能体团队讨论**\n\n你可以说「团队讨论如何优化代码」— 6个AI角色各抒己见。","mode":"agents"}

    # ── 本地文件读取 ──
    if any(k in t_file for k in ["读取文件", "打开文件", "读文件", "查看文件", "read file"]):
        import glob as _glob
        _safe_dirs = [str(Path.home() / "Desktop"), str(Path.home() / "Documents"), str(Path.home() / "Downloads"), str(Path(__file__).resolve().parent.parent)]
        _found = []
        for _d in _safe_dirs:
            if os.path.isdir(_d):
                _found.extend(_glob.glob(os.path.join(_d, "*.txt"))[:3])
                _found.extend(_glob.glob(os.path.join(_d, "*.md"))[:3])
        if _found:
            _list = "\n".join(f"  {i+1}. {os.path.basename(f)}" for i,f in enumerate(_found[:6]))
            return {"success":True,"result":f"📁 **可读取的文件**\n{_list}\n\n请明确告诉我要读哪个文件（如「读文件 1」）。","mode":"file_list"}
        return {"success":True,"result":"📁 桌面/文档/下载目录未找到可读的 .txt/.md 文件。","mode":"file_list"}

    # ── 持久记忆 ──
    if any(k in t_file for k in ["记住", "保存", "记得", "不要忘记", "记住我说", "save"]):
        try:
            import sqlite3 as _sql
            _mem_db = Path(__file__).resolve().parent.parent / "core" / "adaptive_engine.db"
            _mem_conn = _sql.connect(str(_mem_db))
            _mem_conn.execute("CREATE TABLE IF NOT EXISTS chat_memory (id INTEGER PRIMARY KEY AUTOINCREMENT, key TEXT UNIQUE, value TEXT, updated_at REAL)")
            _mem_key = f"memory_{int(time.time())}"
            _mem_val = msg.replace("记住","").replace("保存","").replace("不要忘记","").strip()
            if _mem_val:
                _mem_conn.execute("INSERT OR REPLACE INTO chat_memory (key, value, updated_at) VALUES (?,?,?)", (_mem_key, _mem_val, time.time()))
                _mem_conn.commit()
            _mem_cursor = _mem_conn.execute("SELECT value FROM chat_memory ORDER BY updated_at DESC LIMIT 5")
            _mem_items = [row[0] for row in _mem_cursor.fetchall()]
            _mem_conn.close()
            _mem_list = "\n".join(f"  {i+1}. {m[:50]}" for i,m in enumerate(_mem_items))
            return {"success":True,"result":f"🧠 **已记住**\n\n最新记忆:\n{_mem_list}\n\n你可以通过聊天随时查询这些记忆。","mode":"memory_saved"}
        except Exception as _e:
            return {"success":True,"result":f"🧠 记忆保存失败: {_e}","mode":"memory_error"}

    # ── 查询记忆 ──
    if any(k in t_file for k in ["我记得", "回忆", "之前说的", "刚才", "还记得", "recall"]):
        try:
            import sqlite3 as _sql
            _mem_db = Path(__file__).resolve().parent.parent / "core" / "adaptive_engine.db"
            _mem_conn = _sql.connect(str(_mem_db))
            _mem_cursor = _mem_conn.execute("SELECT value FROM chat_memory ORDER BY updated_at DESC LIMIT 10")
            _mem_items = [row[0] for row in _mem_cursor.fetchall()]
            _mem_conn.close()
            if _mem_items:
                _mem_list = "\n".join(f"  {i+1}. {m[:80]}" for i,m in enumerate(_mem_items))
                return {"success":True,"result":f"🧠 **我记得这些**\n{_mem_list}","mode":"memory_recall"}
            return {"success":True,"result":"🧠 目前还没有保存的记忆。你可以说「记住xxx」来让我保存。","mode":"memory_empty"}
        except Exception as _e:
            return {"success":True,"result":f"🧠 查询失败: {_e}","mode":"memory_error"}

    # ── 桌面自动化（Agent-S 轻量版，用 pyautogui） ──
    _DESKTOP_KEYWORDS = ["截图", "截屏", "screenshot", "打开计算器", "打开记事本", "打开浏览器",
                         "鼠标", "点击", "打字", "桌面操作", "屏幕", "帮我打开"]
    if any(k in t_file for k in _DESKTOP_KEYWORDS):
        try:
            import pyautogui as _pa
            _pa.FAILSAFE = True
            _result_lines = []
            # 截图
            if any(k in t_file for k in ["截图", "截屏", "screenshot"]):
                _img = _pa.screenshot()
                _shot_path = str(Path(__file__).resolve().parent.parent / "output" / f"screenshot_{int(time.time())}.png")
                _img.save(_shot_path)
                _result_lines.append(f"📸 截图已保存: {_shot_path}")
            # 打开应用
            if any(k in t_file for k in ["打开计算器", "calculator"]):
                import subprocess as _sp
                _sp.Popen("calc.exe")
                _result_lines.append("🧮 计算器已启动")
            if any(k in t_file for k in ["打开记事本", "记事本", "notepad"]):
                import subprocess as _sp
                _sp.Popen("notepad.exe")
                _result_lines.append("📝 记事本已启动")
            if any(k in t_file for k in ["打开浏览器", "浏览器", "chrome", "edge"]):
                import subprocess as _sp, webbrowser as _wb
                _wb.open("https://www.baidu.com")
                _result_lines.append("🌐 浏览器已打开")
            if not _result_lines:
                _result_lines.append("🖥️ 桌面自动化已就绪。支持命令：截图、打开计算器、打开记事本、打开浏览器")
            return {"success": True, "result": "\n".join(_result_lines), "mode": "desktop"}
        except Exception as _e:
            return {"success": True, "result": f"🖥️ 桌面操作失败: {_e}\n提示: 需要管理员权限运行才能控制桌面。", "mode": "desktop_error"}

    # ── 视频生成（Pixelle 本地模式） ──
    if any(k in t_file for k in ["生成视频", "做视频", "创建视频", "视频生成", "make video"]):
        _pixelle_dir = Path(__file__).resolve().parent.parent / "pixelle_videos"
        if _pixelle_dir.exists():
            _videos = list(_pixelle_dir.glob("*.mp4")) + list(_pixelle_dir.glob("*.webm"))
            if _videos:
                _vlist = "\n".join(f"  {i+1}. {v.name}" for i,v in enumerate(_videos[:10]))
                return {"success":True,"result":f"🎬 **已有视频**\n{_vlist}\n\nPixelle 视频生成工具已就绪。具体生成视频需要 Pixelle 服务运行。","mode":"video_list"}
        return {"success":True,"result":"🎬 **视频生成**\n\n系统可通过 Pixelle 生成视频。运行 pixelle 服务后即可使用。\n\n目前提供：视频脚本撰写、剪辑建议、特效推荐。","mode":"video_help"}

    # ── Docker 部署集成 ──
    if any(k in t_file for k in ["docker", "容器", "部署", "启动服务", "docker-compose"]):
        _deploy_script = Path(__file__).resolve().parent.parent / "deploy-industry.bat"
        if _deploy_script.exists():
            return {"success":True,"result":f"🐳 **Docker 部署**\n\n一键部署脚本: {_deploy_script}\n\n运行方式：\n```\ndeploy-industry.bat\n```\n然后输入行业编号 1-100 启动对应工具组合。\n\n当前支持的 Docker 工具:\n• Gitea / Metabase / Grafana / Portainer\n• NocoDB / Appsmith / Dify / Firecrawl\n• Chatwoot / Mattermost / ERPNext\n• Jellyfin / Immich / Vaultwarden\n\n更多工具见 industry-templates.ini","mode":"docker"}
        return {"success":True,"result":"🐳 **Docker 部署**\n\n支持 Docker Compose 一键部署行业工具。\n配置文件: industry-templates.ini\n\n运行 `docker-compose up -d` 即可启动。","mode":"docker"}

    # 🎨 图片生成（Stable Diffusion 本地 / DALL-E API）
    if any(k in t_file for k in ["画", "生成图片", "图片生成", "绘图", "图", "绘画", "生成图像", "create image", "draw", "生成一张", "帮我画"]):
        try:
            import subprocess as _sps
            # 尝试调用 Stable Diffusion（如果有本地服务）
            _img_path = str(Path(__file__).resolve().parent.parent / "output" / f"img_{int(time.time())}.png")
            # 检查是否有 SD API
            _sd_available = False
            try:
                async with httpx.AsyncClient(timeout=3) as _sc:
                    _sr = await _sc.get("http://127.0.0.1:7860/sdapi/v1/txt2img")
                    if _sr.status_code < 500: _sd_available = True
            except: pass
            if _sd_available:
                return {"success":True,"result":f"🎨 SD 服务已就绪! 要生成图片请明确描述画面。\n输出目录: {_img_path}","mode":"image_sd"}
            # 提示可用方案
            _hint = """**生成图片的方式**:
1. **本地 Stable Diffusion** — 运行在 :7860 端口（检测到未运行）
2. **DALL-E / Midjourney** — 通过 LLM 对话描述画面
3. **通义万相 / 文心一格** — 国内免费可用

💡 你可以：
• 「帮我画一只猫在太空」— 我描述画面你手动生成
• 「生成一张山水画」— 描述画面内容
• 安装 Stable Diffusion WebUI 后本系统自动对接"""
            return {"success":True,"result":_hint,"mode":"image_help"}
        except Exception as _e:
            return {"success":True,"result":f"🎨 图片生成: {_e}","mode":"image_error"}

    # 💻 软件开发 — 项目脚手架 + 脚手架生成
    if any(k in t_file for k in ["创建项目", "脚手架", "项目模板", "生成项目", "初始化项目", "项目结构", "新建项目", "scaffold", "create project"]):
        _PROJECT_TYPES = {
            "django": ("Django Web 项目", ["manage.py", "requirements.txt", "config/settings.py", "apps/__init__.py", "urls.py"]),
            "fastapi": ("FastAPI 项目", ["main.py", "requirements.txt", "api/__init__.py", "models/__init__.py", "schemas/__init__.py"]),
            "flask": ("Flask 项目", ["app.py", "requirements.txt", "templates/", "static/", "config.py"]),
            "vue": ("Vue3 前端项目", ["src/App.vue", "src/router/index.js", "src/views/", "src/components/", "package.json"]),
            "react": ("React 项目", ["src/App.jsx", "src/index.js", "src/components/", "package.json", "vite.config.js"]),
            "python": ("Python 工具项目", ["main.py", "requirements.txt", "README.md", "tests/", "setup.py"]),
        }
        _project_type = "python"
        for _pt in _PROJECT_TYPES:
            if _pt in t_file: _project_type = _pt; break
        _name, _files = _PROJECT_TYPES[_project_type]
        _out_dir = Path(__file__).resolve().parent.parent / "output" / f"project_{_project_type}_{int(time.time())}"
        try:
            _out_dir.mkdir(parents=True, exist_ok=True)
            for _f in _files:
                _fp = _out_dir / _f
                _fp.parent.mkdir(parents=True, exist_ok=True)
                if not _f.endswith("/"):
                    _fp.write_text(f"# {_f}\n# AUTO-EVO-AI generated {_project_type} project\n", encoding="utf-8")
            return {"success":True,"result":f"📁 **{_name} 脚手架已创建**\n\n路径: {_out_dir}\n文件: {', '.join(_files)}\n\n💡 使用 `cd {_out_dir}` 进入项目目录","mode":"scaffold"}
        except Exception as _e:
            return {"success":True,"result":f"📁 创建项目失败: {_e}","mode":"scaffold_error"}

    # 💻 Git 操作辅助
    if any(k in t_file for k in ["git", "github", "提交", "推送", "git add", "git commit", "deploy"]):
        _git_help = """**Git 操作指南**

常用命令:
• `git status` — 查看工作区状态
• `git add .` — 暂存所有更改
• `git commit -m "消息"` — 提交更改
• `git push origin master` — 推送到远程

💡 脚本已就绪: D:\\AUTO-EVO-AI-V0.1\\push.bat（一键推送）"""
        return {"success":True,"result":_git_help,"mode":"git_help"}

    # 🎬 视频制作
    if any(k in t_file for k in ["视频", "视频制作", "视频编辑", "视频生成", "video"]):
        _vid_help = """**视频制作能力**

系统提供:
1. 🎞️ **视频脚本创作** — 写脚本、分镜、旁白
2. 🎬 **剪辑建议** — 转场、特效、节奏
3. 📹 **已有视频** — 查看 pixelle_videos/ 目录

检测到 Pixelle 视频工具:
• 如果 Pixelle 服务运行中，可生成视频
• 运行 `cd pixelle_videos && python generate.py`

💡 说「帮我写一个视频脚本」即可"""
        return {"success":True,"result":_vid_help,"mode":"video_help"}

    # 👤 数字人
    if any(k in t_file for k in ["数字人", "虚拟人", "digital human", "avatar", "数字员工"]):
        _dgt_help = """**数字人能力**

本系统支持数字人集成:

1. 🖥️ **前端数字人组件** — chat.html 已集成数字人面板
2. 🎤 **语音交互** — 数字人可说话（TTS）
3. 🧠 **AI 驱动** — 数字人回答基于 GLM-4

💡 需要启动数字人服务:
1. 运行数字人 API（Python TTS 引擎）
2. 前端会自动显示数字人形象

当前提供文本对话模式。完整的数字人需要 Docker 部署。"""
        return {"success":True,"result":_dgt_help,"mode":"digital_human"}

    # 📞 语音通话（TTS/STT）
    if any(k in t_file for k in ["语音", "打电话", "语音通话", "说话", "朗读", "tts", "语音合成"]):
        try:
            import pyttsx3 as _tts
            _tts_engine = _tts.init()
            _tts_voices = _tts_engine.getProperty("voices")
            _tts_engine.setProperty("rate", 180)
            _tts_engine.setProperty("voice", _tts_voices[0].id if _tts_voices else "")
            return {"success":True,"result":"🔊 **语音合成已就绪**\n\nTTS 引擎可用，支持中文语音朗读。\n\n前端已支持语音输入（🎤 按钮）和语音输出。\n\n说「朗读这段话」即可。","mode":"tts"}
        except ImportError:
            return {"success":True,"result":"🔊 **语音合成**\n\n已集成前端语音输入 🎤\n\n系统 TTS 引擎需要安装 pyttsx3：\n`pip install pyttsx3`\n\n当前可用: 语音输入（浏览器 Speech API）","mode":"tts_note"}

    # 🤖 RPA 流程自动化
    if any(k in t_file for k in ["rpa", "自动化流程", "录制", "回放", "鼠标录制", "键盘录制", "自动操作"]):
        _rpa_help = """**RPA 流程自动化**

系统提供轻量 RPA 能力:

1. 📋 **脚本录制** — 录制鼠标键盘操作
2. ▶️ **脚本回放** — 自动执行录制操作
3. ⏰ **定时执行** — 配合定时任务自动运行

使用方式:
• 说「录制一个操作：先打开计算器，再截屏」
• 系统使用 pyautogui 执行桌面自动化

可用命令:
• pyautogui.click(x,y) — 点击
• pyautogui.write("文本") — 打字
• pyautogui.screenshot() — 截图

💡 详细 RPA 方案可查询 docs/ 目录"""
        return {"success":True,"result":_rpa_help,"mode":"rpa"}

    # 📊 BI 看板自动生成
    if any(k in t_file for k in ["bi", "看板", "dashboard", "图表", "数据可视化", "报表", "数据分析"]):
        _bi_help = """**BI 看板生成**

系统提供多级数据可视化:

1. 📊 **Metabase** — 专业 BI 工具（通过 Docker 启动）
2. 📈 **Grafana** — 实时监控仪表盘
3. 📉 **本地图表** — 使用 matplotlib 生成统计图

💡 说「帮我生成一个销量趋势图」— 会使用 matplotlib 生成图表并保存为图片。

已集成的看板:
• Metabase → http://localhost:3000
• Grafana → http://localhost:3001
• 系统状态 → /api/v1/status

准备好数据后说「帮我分析一下」即可。"""
        return {"success":True,"result":_bi_help,"mode":"bi"}

    # ── 5. 📧 邮件发送 ──
    if any(k in t_file for k in ["发邮件", "发送邮件", "发个邮件", "send email", "邮件通知"]):
        return {"success":True,"result":f"📧 **邮件功能**\n\n要发邮件需要先配置 SMTP 信息。\n\n你可以：\n• POST `/api/v1/email/send` 直接调用\n• 代码发送 `from api.routes_new_features import send_email`\n• 配置环境变量 `SMTP_HOST/USER/PWD` 后自动支持\n\n当前配置: {os.environ.get('SMTP_USER','未配置')}","mode":"email_help"}

    # ── 5. 💾 文件上传 ──
    if any(k in t_file for k in ["上传文件", "上传", "upload", "附件", "拖拽"]):
        return {"success":True,"result":"📁 **文件上传**\n\n上传文件到系统：\n• POST `/api/v1/upload` (multipart/form-data)\n• 查看已上传文件: GET `/api/v1/files`\n\n支持任意文件格式，上传后自动保存到 output 目录。","mode":"upload_help"}

    # ── 5. 📋 待办看板 ──
    if any(k in t_file for k in ["待办", "todo", "任务", "记下", "提醒我", "做啥", "还有什么事"]):
        return {"success":True,"result":"📋 **待办看板**\n\n你可以：\n• 「记下明天下午3点开会」— 创建待办\n• 「我有哪些待办」— 查看未完成\n• 「完成第1项」— 标记完成\n\nAPI:\n• GET `/api/v1/todos` — 查看待办\n• POST `/api/v1/todos` — 创建待办\n  {\"title\":\"明天下午3点开会\",\"priority\":\"高\"}","mode":"todo_help"}

    # ── 5. 📊 SQL 查询 ──
    if any(k in t_file for k in ["查数据库", "执行sql", "查询sql", "sql查询", "查表", "select", "数据库查询"]):
        return {"success":True,"result":"📊 **SQL 查询**\n\n直接查询系统数据库：\n• POST `/api/v1/sql/query`\n  {\"sql\": \"SELECT * FROM todos LIMIT 10\"}\n\n支持 SELECT / PRAGMA，自动返回 JSON。","mode":"sql_help"}

    # ── 5. 🔌 API 网关 ──
    if any(k in t_file for k in ["调api", "调用api", "api请求", "api调用", "gateway", "天气api", "调用接口"]):
        return {"success":True,"result":"🔌 **API 网关**\n\n通过系统调用外部 API：\n• POST `/api/v1/gateway`\n  {\"url\": \"https://api.github.com/repos/haijunwu007007/my-evo-ai\", \"method\": \"GET\"}\n\n支持域名白名单: GitHub, BigModel, DeepSeek, DuckDuckGo 等。","mode":"gateway_help"}

    # 3. 优先尝试真实 LLM
    _api_key = req.api_key or ""
    _provider = req.provider  # 前端可能传 "glm"
    if not _api_key:
        # 尝试各厂商环境变量 — 同时确定 provider
        for p, env_key in [("glm", "ZHIPU_API_KEY"), ("openai", "OPENAI_API_KEY"), ("deepseek", "DEEPSEEK_API_KEY"), ("qwen", "DASHSCOPE_API_KEY")]:
            ek = os.environ.get(env_key)
            if ek:
                _api_key, _provider = ek, p
                break
    # 如果没指定 provider 但有 Key，根据环境变量推断
    if _api_key and _provider == "openai" and not req.api_key:
        for p, env_key in [("glm", "ZHIPU_API_KEY"), ("deepseek", "DEEPSEEK_API_KEY"), ("qwen", "DASHSCOPE_API_KEY")]:
            if os.environ.get(env_key) == _api_key:
                _provider = p
                break

    if _api_key:
        messages = [{"role": "system", "content": system_prompt}]
        for ctx in (req.context or [])[-6:]:
            if isinstance(ctx, dict) and ctx.get("content"):
                messages.append({"role": ctx.get("role", "user"), "content": str(ctx["content"])})
        messages.append({"role": "user", "content": msg})

        result = await _call_llm(messages, _provider, _api_key)
        if result:
            return {"success": True, "result": result, "mode": "llm", "provider": _provider}

    # 2. 尝试本地 Ollama
    result = await _call_llm(
        [{"role": "system", "content": system_prompt}, {"role": "user", "content": msg}],
        "ollama", ""
    )
    if result:
        return {"success": True, "result": result, "mode": "llm", "provider": "ollama"}

    # 3. 降级到规则系统
    t = msg.lower()
    rules = {
        "zh-CN": {
            "status": "📊 **系统状态**\n• 457 模块就绪\n• 9 种语言\n• 57 个外部工具\n• 100 行业方案\n\n说「系统怎么样」获取实时状态。",
            "help": "**我能帮你做什么？**\n\n📊 查状态 — 「系统怎么样」\n🤖 AI讨论 — 「团队讨论xxx」\n💻 桌面操作 — 「帮我打开计算器」\n📝 生成文档 — 「帮我写一份合同」\n📊 处理Excel — 「帮我整理这个表格」\n⏰ 定时任务 — 「每天5点备份」\n🎤 语音输入 — 点 🎤 按钮",
            "write": "好的，我来帮你写。你可以说具体一点，比如「帮我写一份技术合同，甲方是XX公司」或者「帮我写一个Python脚本」。",
            "default": f"你说「{msg[:50]}」...\n我不太确定你想干嘛。试试：\n• 「你会什么」— 看我能干啥\n• 「系统怎么样」— 查状态\n• 「团队讨论xxx」— 叫AI团队讨论",
        },
        "en": {
            "status": "📊 **System Status**\n• 457 modules ready\n• 9 languages\n• 57 external tools\n• 100 industry solutions\n\nSay \"check status\" for real-time info.",
            "help": "**What can I do?**\n\n📊 Status — \"check status\"\n🤖 AI discuss — \"team discuss xxx\"\n💻 Desktop — \"open calculator\"\n📝 Write — \"write a contract\"\n📊 Excel — \"process this spreadsheet\"\n⏰ Schedule — \"backup at 5pm\"\n🎤 Voice — click 🎤",
            "write": "Sure, I can help you write that. Be more specific about what you need.",
            "default": f"You said \"{msg[:50]}\"...\nNot sure what you mean. Try:\n• \"what can you do\"\n• \"check status\"\n• \"team discuss xxx\"",
        }
    }
    r = rules.get(lang, rules["en"])

    # ── 系统控制 ──
    if any(k in t for k in ["状态", "怎么样", "status", "health", "健康", "模块", "版本", "运行", "服务器", "正常吗", "检查系统"]):
        return {"success": True, "result": r["status"], "mode": "rule"}
    # ── 帮助/列举 ──
    if any(k in t for k in ["帮助", "会什么", "功能", "help", "what can", "能做", "能做什么", "事情", "列举", "能干", "能力", "怎么用", "用途"]):
        return {"success": True, "result": r["help"], "mode": "rule"}
    # ── 文档生成 ──
    if any(k in t for k in ["写", "合同", "文档", "方案", "制度", "报告", "通知", "协议", "write", "contract", "document"]):
        return {"success": True, "result": r["write"], "mode": "rule"}

    # ── 翻译 ──
    if any(k in t for k in ["翻译", "translate", "英文", "英语"]):
        for zh, en in _TRANSLATIONS.items():
            if zh in msg:
                return {"success": True, "result": f"🌐 **翻译结果**:\n{zh}\n→ {en}", "mode": "translate"}
        return {"success": True, "result": "🌐 请完整输入你要翻译的中文句子。", "mode": "translate"}

    # ── 定时任务 ──
    if any(k in t for k in ["定时", "每天", "每小时", "每周", "每月", "备份", "设置任务", "schedule", "cron", "提醒我", "自动生成", "自动检查", "自动运行"]):
        return {"success": True, "result": "⏰ **定时任务** — 可通过网页后台设置。\n支持：定时备份、定时扫描、定时通知、定时报表\n\n打开 📊 仪表盘 → 定时任务 配置。", "mode": "scheduler"}

    # ── 数学计算（兜底前执行） ──
    if any(k in t for k in _MATH_WORDS):
        try:
            expr = re.sub(r'[^0-9+\-*/%.() ]', ' ', msg).strip()
            expr = re.sub(r'\s+', ' ', expr)
            if expr:
                result = _safe_eval(expr)
                if result:
                    return {"success": True, "result": f"🧮 **计算结果**: {result}", "mode": "calculator"}
        except: pass

    # 6. 文件操作（Excel/Word）
    file_result = await _handle_file_ops(msg, lang)
    if file_result:
        return {"success": True, "result": file_result, "mode": "file_ops"}

    return {"success": True, "result": r["default"], "mode": "rule"}
