"""
AUTO-EVO-AI V0.1 — 一键PPT生成
对话式输入主题 → AI搜索 → LLM规划大纲 → 写内容 → python-pptx生成 → 下载
"""
from fastapi import APIRouter
from pydantic import BaseModel
import os, json, uuid, subprocess, re
from pathlib import Path
from core.logging_config import get_logger

logger = get_logger("evo.api.ppt_gen")
router = APIRouter(prefix="/api/v1/ppt-auto", tags=["ppt_auto"])

BASE = Path(__file__).resolve().parent.parent.parent
OUTPUT_DIR = BASE / "output" / "ppt"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

class GenerateRequest(BaseModel):
    topic: str

@router.post("/generate")
def generate_ppt(req: GenerateRequest):
    topic = req.topic.strip()
    if not topic:
        return {"success": False, "error": "请输入PPT主题"}
    
    try:
        # 1. 生成大纲
        outline = _generate_outline(topic)
        
        # 2. 写每页内容
        slides = _generate_slides(topic, outline)
        
        # 3. 生成 PPTX
        filename, ppt_path = _create_pptx(topic, slides)
        
        word_count = sum(len(s.get("content","")) for s in slides)
        
        return {
            "success": True,
            "filename": filename,
            "outline": [{"title": s["title"], "desc": s.get("desc","")} for s in slides],
            "slides_count": len(slides),
            "word_count": word_count,
            "download_url": f"/output/ppt/{filename}"
        }
    except Exception as e:
        logger.error(f"[PPT] generate error: {e}")
        return {"success": False, "error": str(e)[:200]}

def _generate_outline(topic: str) -> list:
    """LLM 生成 PPT 大纲，超时降级模板"""
    try:
        from api.agent_llm import call_llm
        import threading, queue
        q = queue.Queue()
        def _do():
            try:
                text, _ = call_llm([
                    {"role": "system", "content": "你是专业的PPT策划专家，输出结构化的JSON大纲。"},
                    {"role": "user", "content": f"为「{topic}」生成PPT大纲。8-15页，每页标题+描述。JSON格式：[{{'title':'xxx','desc':'xxx'}}]"}
                ], None, "")
                q.put(text or '[]')
            except: q.put('')
        t = threading.Thread(target=_do, daemon=True)
        t.start()
        t.join(timeout=15)
        if not t.is_alive():
            text = q.get_nowait() if not q.empty() else ''
            json_match = re.search(r'\[.*\]', text, re.DOTALL)
            if json_match:
                outline = json.loads(json_match.group())
                if isinstance(outline, list) and len(outline) >= 3:
                    return outline
    except: pass
    
    # 降级模板
    return [
        {"title": f"{topic} 概述", "desc": "背景介绍与核心观点"},
        {"title": "市场现状", "desc": "当前市场规模与主要参与者"},
        {"title": "趋势分析", "desc": "关键趋势与驱动因素"},
        {"title": "挑战与机遇", "desc": "面临的主要挑战与潜在机会"},
        {"title": "结论与建议", "desc": "总结核心观点与行动建议"},
    ]

def _generate_slides(topic: str, outline: list) -> list:
    """为每页写内容"""
    from api.agent_llm import call_llm
    slides = []
    for i, item in enumerate(outline):
        prompt = f"""为PPT第{i+1}页写内容。主题：{topic}，标题：{item['title']}
要求：
- 100-200字中文正文
- 3-5个要点（每个要点10-30字）
- 直接返回JSON：{{"title":"{item['title']}","desc":"{item['desc']}","content":"正文...","bullets":["要点1","要点2"]}}"""
        
        text, _ = call_llm([
            {"role": "system", "content": "你是PPT内容撰写专家，输出JSON。"},
            {"role": "user", "content": prompt}
        ], None, "")
        
        try:
            js = json.loads(re.search(r'\{.*\}', text or '{}', re.DOTALL).group())
            slides.append(js)
        except:
            slides.append({
                "title": item["title"],
                "desc": item.get("desc",""),
                "content": f"关于{topic}的{item['title']}分析。",
                "bullets": ["关键数据", "核心发现", "趋势分析"]
            })
    return slides

def _create_pptx(topic: str, slides: list) -> tuple:
    """生成 PPTX 文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = {
        "bg": RGBColor(0x0F, 0x0F, 0x1A),
        "card": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x43, 0x61, 0xEE),
        "text": RGBColor(0xE0, 0xE0, 0xE0),
        "text2": RGBColor(0x88, 0x92, 0xB0),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
    }
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colors["bg"]
    
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(40)
    p.font.color.rgb = colors["white"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AUTO-EVO-AI · 一键生成"
    p2.font.size = Pt(18)
    p2.font.color.rgb = colors["accent"]
    p2.alignment = PP_ALIGN.CENTER
    
    # 内容页
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = colors["bg"]
        
        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = s.get("title", "")
        p.font.size = Pt(28)
        p.font.color.rgb = colors["accent"]
        p.font.bold = True
        
        # 正文
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(2))
        tf2 = txBox2.text_frame
        p = tf2.paragraphs[0]
        p.text = s.get("content", "")
        p.font.size = Pt(16)
        p.font.color.rgb = colors["text"]
        p.space_after = Pt(12)
        
        # 要点
        for b in s.get("bullets", []):
            p = tf2.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(14)
            p.font.color.rgb = colors["text2"]
            p.space_after = Pt(6)
        
        # 页脚
        txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5))
        tf3 = txBox3.text_frame
        p = tf3.paragraphs[0]
        p.text = f"AUTO-EVO-AI · {topic}"
        p.font.size = Pt(10)
        p.font.color.rgb = colors["text2"]
    
    fid = uuid.uuid4().hex[:12]
    filename = f"ppt_{fid}.pptx"
    ppt_path = str(OUTPUT_DIR / filename)
    prs.save(ppt_path)
    return filename, ppt_path
